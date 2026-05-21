"""Document processor — load, parse, clean, and chunk documents of multiple types."""

import hashlib
import os
import re
import unicodedata
from datetime import datetime
from pathlib import Path
from typing import Optional
from urllib.parse import urlparse

import yaml
from loguru import logger
from tqdm import tqdm

from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".txt", ".csv", ".pptx", ".md", ".html"}


class DocumentProcessor:
    """Load, parse, clean, and chunk documents from various sources.

    Supports PDF, DOCX, PPTX, TXT, CSV, HTML, and web URLs.
    """

    def __init__(self, config: dict) -> None:
        """Initialise with config dict (from config.yaml)."""
        self.config = config
        chunking = config.get("chunking", {})
        self.chunk_size: int = chunking.get("chunk_size", 1000)
        self.chunk_overlap: int = chunking.get("chunk_overlap", 200)
        self.separators: list[str] = chunking.get(
            "separators", ["\n\n", "\n", ".", "!", "?", ",", " "]
        )
        self.splitter = RecursiveCharacterTextSplitter(
            chunk_size=self.chunk_size,
            chunk_overlap=self.chunk_overlap,
            separators=self.separators,
            length_function=len,
        )

    # ─────────────────────────────────────────────────────────────────────────
    # Public API
    # ─────────────────────────────────────────────────────────────────────────

    def load_document(self, file_path: str) -> list[Document]:
        """Auto-detect file type and load document with metadata.

        Args:
            file_path: Absolute or relative path to the document file.

        Returns:
            List of LangChain Document objects (one per page for PDFs).
        """
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        ext = path.suffix.lower()
        logger.info(f"Loading {ext} document: {path.name}")

        dispatchers = {
            ".pdf": self.extract_pdf,
            ".docx": self.extract_docx,
            ".pptx": self.extract_pptx,
            ".txt": self._extract_txt,
            ".md": self._extract_txt,
            ".csv": self._extract_csv,
            ".html": self._extract_html_file,
        }

        if ext not in dispatchers:
            raise ValueError(f"Unsupported file type: {ext}")

        docs = dispatchers[ext](str(path))
        logger.info(f"Loaded {len(docs)} page(s)/section(s) from {path.name}")
        return docs

    def load_multiple_documents(self, directory: str) -> list[Document]:
        """Load all supported documents from a directory.

        Args:
            directory: Path to a directory containing documents.

        Returns:
            Flat list of all loaded Document objects.
        """
        dir_path = Path(directory)
        if not dir_path.is_dir():
            raise NotADirectoryError(f"Not a directory: {directory}")

        files = [
            p for p in dir_path.rglob("*")
            if p.is_file() and p.suffix.lower() in SUPPORTED_EXTENSIONS
        ]
        logger.info(f"Found {len(files)} supported files in {directory}")

        all_docs: list[Document] = []
        errors: list[str] = []

        for file_path in tqdm(files, desc="Loading documents"):
            try:
                docs = self.load_document(str(file_path))
                all_docs.extend(docs)
            except Exception as exc:
                logger.warning(f"Failed to load {file_path.name}: {exc}")
                errors.append(str(file_path))

        if errors:
            logger.warning(f"Failed to load {len(errors)} file(s): {errors}")

        logger.info(f"Successfully loaded {len(all_docs)} document sections total.")
        return all_docs

    def extract_pdf(self, file_path: str) -> list[Document]:
        """Extract text page-by-page from a PDF.

        Args:
            file_path: Path to the PDF file.

        Returns:
            List of Documents, one per page.
        """
        try:
            import fitz  # PyMuPDF
        except ImportError:
            raise ImportError("PyMuPDF not installed. Run: pip install pymupdf")

        path = Path(file_path)
        docs: list[Document] = []

        with fitz.open(file_path) as pdf:
            for page_num in range(len(pdf)):
                page = pdf[page_num]
                text = page.get_text("text")

                if not text.strip():
                    logger.debug(f"Page {page_num + 1} in {path.name} has no text (possibly scanned).")
                    continue

                text = self.clean_text(text)
                meta = self._build_metadata(
                    path,
                    extra={
                        "page": page_num + 1,
                        "total_pages": len(pdf),
                    },
                )
                docs.append(Document(page_content=text, metadata=meta))

        return docs

    def extract_docx(self, file_path: str) -> list[Document]:
        """Extract paragraphs and tables from a DOCX file.

        Args:
            file_path: Path to the DOCX file.

        Returns:
            List of Documents (one per logical section).
        """
        try:
            import docx
        except ImportError:
            raise ImportError("python-docx not installed. Run: pip install python-docx")

        path = Path(file_path)
        doc = docx.Document(file_path)
        sections: list[str] = []
        current: list[str] = []

        for para in doc.paragraphs:
            style = para.style.name.lower() if para.style else ""
            text = para.text.strip()
            if not text:
                continue
            if "heading" in style and current:
                sections.append("\n".join(current))
                current = [text]
            else:
                current.append(text)

        if current:
            sections.append("\n".join(current))

        meta = self._build_metadata(path)
        return [
            Document(
                page_content=self.clean_text(sec),
                metadata={**meta, "section": i + 1},
            )
            for i, sec in enumerate(sections)
            if sec.strip()
        ]

    def extract_pptx(self, file_path: str) -> list[Document]:
        """Extract text from each slide of a PPTX file."""
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("python-pptx not installed. Run: pip install python-pptx")

        path = Path(file_path)
        prs = Presentation(file_path)
        docs: list[Document] = []
        meta = self._build_metadata(path)

        for slide_num, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                content = self.clean_text("\n".join(texts))
                docs.append(
                    Document(
                        page_content=content,
                        metadata={**meta, "slide": slide_num},
                    )
                )
        return docs

    def extract_url(self, url: str) -> Document:
        """Fetch and extract text content from a web URL.

        Args:
            url: The URL to fetch.

        Returns:
            A single Document with cleaned web content.
        """
        try:
            import requests
            from bs4 import BeautifulSoup
        except ImportError:
            raise ImportError(
                "requests and beautifulsoup4 required. "
                "Run: pip install requests beautifulsoup4"
            )

        logger.info(f"Fetching URL: {url}")
        resp = requests.get(url, timeout=15, headers={"User-Agent": "DocuMindAI/1.0"})
        resp.raise_for_status()

        soup = BeautifulSoup(resp.text, "html.parser")
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()

        text = soup.get_text(separator="\n", strip=True)
        title = soup.title.string if soup.title else urlparse(url).netloc

        return Document(
            page_content=self.clean_text(text),
            metadata={
                "source": url,
                "title": title,
                "file_type": "url",
                "loaded_at": datetime.now().isoformat(),
            },
        )

    def chunk_documents(self, documents: list[Document]) -> list[Document]:
        """Split documents into smaller chunks for embedding.

        Args:
            documents: List of loaded Document objects.

        Returns:
            List of chunked Document objects with chunk metadata.
        """
        logger.info(f"Chunking {len(documents)} document(s) ...")
        chunks = self.splitter.split_documents(documents)

        for idx, chunk in enumerate(chunks):
            chunk.metadata["chunk_id"] = f"chunk_{idx:05d}"
            chunk.metadata["chunk_index"] = idx
            chunk.metadata["document_hash"] = _hash_text(chunk.page_content)

        logger.info(f"Created {len(chunks)} chunks (size={self.chunk_size}, overlap={self.chunk_overlap}).")
        return chunks

    def get_document_stats(self, documents: list[Document]) -> dict:
        """Compute statistics over a list of documents.

        Args:
            documents: List of Document objects.

        Returns:
            Dict with total_docs, total_chunks, word counts, file types, etc.
        """
        if not documents:
            return {}

        word_counts = [len(d.page_content.split()) for d in documents]
        file_types: dict[str, int] = {}
        sources: list[str] = []

        for doc in documents:
            ft = doc.metadata.get("file_type", "unknown")
            file_types[ft] = file_types.get(ft, 0) + 1
            src = doc.metadata.get("source", "")
            if src and src not in sources:
                sources.append(src)

        return {
            "total_docs": len(sources),
            "total_chunks": len(documents),
            "total_words": sum(word_counts),
            "avg_chunk_size": round(sum(len(d.page_content) for d in documents) / len(documents)),
            "avg_words_per_chunk": round(sum(word_counts) / len(word_counts)),
            "file_types": file_types,
            "largest_chunk_words": max(word_counts),
            "smallest_chunk_words": min(word_counts),
        }

    def deduplicate_documents(self, documents: list[Document]) -> list[Document]:
        """Remove duplicate chunks based on content hash.

        Args:
            documents: List of possibly-duplicate Document objects.

        Returns:
            Deduplicated list preserving original order.
        """
        seen: set[str] = set()
        unique: list[Document] = []
        for doc in documents:
            h = _hash_text(doc.page_content)
            if h not in seen:
                seen.add(h)
                unique.append(doc)

        removed = len(documents) - len(unique)
        if removed:
            logger.info(f"Removed {removed} duplicate chunk(s).")
        return unique

    def clean_text(self, text: str) -> str:
        """Normalise and clean raw extracted text.

        Args:
            text: Raw text string.

        Returns:
            Cleaned, normalised text.
        """
        # Unicode normalisation
        text = unicodedata.normalize("NFKC", text)
        # Fix common encoding artefacts
        text = text.replace("\x00", "").replace("﻿", "")
        # Collapse repeated whitespace / blank lines
        text = re.sub(r"[ \t]+", " ", text)
        text = re.sub(r"\n{3,}", "\n\n", text)
        # Strip header/footer noise patterns (page numbers etc.)
        text = re.sub(r"(?m)^\s*\d+\s*$", "", text)
        text = text.strip()
        return text

    # ─────────────────────────────────────────────────────────────────────────
    # Private helpers
    # ─────────────────────────────────────────────────────────────────────────

    def _extract_txt(self, file_path: str) -> list[Document]:
        path = Path(file_path)
        text = path.read_text(encoding="utf-8", errors="replace")
        meta = self._build_metadata(path)
        return [Document(page_content=self.clean_text(text), metadata=meta)]

    def _extract_csv(self, file_path: str) -> list[Document]:
        try:
            import pandas as pd
        except ImportError:
            raise ImportError("pandas not installed. Run: pip install pandas")

        path = Path(file_path)
        df = pd.read_csv(file_path)
        text = df.to_string(index=False)
        meta = self._build_metadata(path, extra={"rows": len(df), "columns": len(df.columns)})
        return [Document(page_content=self.clean_text(text), metadata=meta)]

    def _extract_html_file(self, file_path: str) -> list[Document]:
        try:
            from bs4 import BeautifulSoup
        except ImportError:
            raise ImportError("beautifulsoup4 not installed. Run: pip install beautifulsoup4")

        path = Path(file_path)
        html = path.read_text(encoding="utf-8", errors="replace")
        soup = BeautifulSoup(html, "html.parser")
        for tag in soup(["script", "style"]):
            tag.decompose()
        text = soup.get_text(separator="\n", strip=True)
        meta = self._build_metadata(path)
        return [Document(page_content=self.clean_text(text), metadata=meta)]

    @staticmethod
    def _build_metadata(path: Path, extra: Optional[dict] = None) -> dict:
        stat = path.stat()
        meta = {
            "source": str(path),
            "filename": path.name,
            "file_type": path.suffix.lower().lstrip("."),
            "file_size_bytes": stat.st_size,
            "loaded_at": datetime.now().isoformat(),
        }
        if extra:
            meta.update(extra)
        return meta


# ─────────────────────────────────────────────────────────────────────────────
# Utility
# ─────────────────────────────────────────────────────────────────────────────

def _hash_text(text: str) -> str:
    return hashlib.md5(text.encode("utf-8")).hexdigest()


# ─────────────────────────────────────────────────────────────────────────────
# Config loader helper (used by other modules too)
# ─────────────────────────────────────────────────────────────────────────────

def load_config(config_path: str = "config/config.yaml") -> dict:
    """Load YAML config file."""
    with open(config_path, "r", encoding="utf-8") as f:
        return yaml.safe_load(f)


if __name__ == "__main__":
    cfg = load_config()
    processor = DocumentProcessor(cfg)
    docs = processor.load_multiple_documents("data/raw/sample_docs")
    chunks = processor.chunk_documents(docs)
    chunks = processor.deduplicate_documents(chunks)
    stats = processor.get_document_stats(chunks)
    print(f"\nDocument Stats: {stats}")
